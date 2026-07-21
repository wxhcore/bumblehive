import asyncio

import pytest

from bumblehive.protocols import ToolCall
from bumblehive.tools import PathAllowlist, ToolManager
from bumblehive.tools.builtins.workspace import WorkspaceAccess


def _manager():
    manager = ToolManager()
    manager.register_builtin_tools()
    return manager


async def _execute(manager, workspace, name, arguments, *, allowlist=PathAllowlist()):
    return await manager.execute_call(
        ToolCall(f"call-{name}", name, arguments),
        workspace=workspace,
        path_allowlist=allowlist,
    )


@pytest.mark.asyncio
async def test_text_file_flow_covers_write_read_dedup_pagination_and_edit(tmp_path) -> None:
    manager = _manager()

    written = await _execute(
        manager,
        tmp_path,
        "write_file",
        {"path": "notes.txt", "content": "one\ntwo\nthree\n"},
    )
    page = await _execute(
        manager,
        tmp_path,
        "read_file",
        {"path": "notes.txt", "offset": 2, "limit": 1},
    )
    duplicate = await _execute(
        manager,
        tmp_path,
        "read_file",
        {"path": "notes.txt", "offset": 2, "limit": 1},
    )
    edited = await _execute(
        manager,
        tmp_path,
        "edit_file",
        {"path": "notes.txt", "old_text": "two", "new_text": "second"},
    )

    assert written.content["success"] is True
    assert page.content["content"] == "2| two"
    assert page.content["next_offset"] == 3
    assert duplicate.content["deduplicated"] is True
    assert edited.content["success"] is True
    assert "warning" not in edited.content
    assert (tmp_path / "notes.txt").read_text(encoding="utf-8") == "one\nsecond\nthree\n"


@pytest.mark.asyncio
async def test_edit_warns_before_read_and_reports_ambiguous_replacements(tmp_path) -> None:
    target = tmp_path / "notes.txt"
    target.write_text("same\nsame\n", encoding="utf-8")
    manager = _manager()

    ambiguous = await _execute(
        manager,
        tmp_path,
        "edit_file",
        {"path": "notes.txt", "old_text": "same", "new_text": "changed"},
    )
    selected = await _execute(
        manager,
        tmp_path,
        "edit_file",
        {
            "path": "notes.txt",
            "old_text": "same",
            "new_text": "changed",
            "occurrence": 2,
        },
    )

    assert ambiguous.content["error"].startswith("old_text appears multiple times")
    assert selected.content["success"] is True
    assert "not been read" in selected.content["warning"]
    assert target.read_text(encoding="utf-8") == "same\nchanged\n"


@pytest.mark.asyncio
async def test_read_file_extracts_pdf_and_office_documents(tmp_path) -> None:
    import fitz
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    from pptx.util import Inches

    pdf_path = tmp_path / "sample.pdf"
    pdf = fitz.open()
    pdf.new_page().insert_text((72, 72), "PDF first")
    pdf.new_page().insert_text((72, 72), "PDF second")
    pdf.save(pdf_path)
    pdf.close()

    document = Document()
    document.add_paragraph("DOCX paragraph")
    document.save(tmp_path / "sample.docx")

    workbook = Workbook()
    workbook.active.title = "Data"
    workbook.active.append(["Name", "Value"])
    workbook.active.append(["bumble", 42])
    workbook.save(tmp_path / "sample.xlsx")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "PPTX text"
    presentation.save(tmp_path / "sample.pptx")

    manager = _manager()
    pdf_result = await _execute(manager, tmp_path, "read_file", {"path": "sample.pdf", "pages": "2"})
    docx_result = await _execute(manager, tmp_path, "read_file", {"path": "sample.docx"})
    xlsx_result = await _execute(manager, tmp_path, "read_file", {"path": "sample.xlsx"})
    pptx_result = await _execute(manager, tmp_path, "read_file", {"path": "sample.pptx"})

    assert pdf_result.content["pages"] == "2-2"
    assert "PDF second" in pdf_result.content["content"]
    assert "PDF first" not in pdf_result.content["content"]
    assert "DOCX paragraph" in docx_result.content["content"]
    assert "bumble\t42" in xlsx_result.content["content"]
    assert "PPTX text" in pptx_result.content["content"]


@pytest.mark.asyncio
async def test_file_tools_use_allowlisted_roots_for_read_list_write_and_edit(tmp_path) -> None:
    workspace = tmp_path / "workspace"
    skills = tmp_path / "skills"
    workspace.mkdir()
    skill_dir = skills / "audit"
    skill_dir.mkdir(parents=True)
    skill_file = skill_dir / "SKILL.md"
    skill_file.write_text("original\n", encoding="utf-8")
    allowlist = PathAllowlist.from_roots(extra_write_roots=[skills])
    manager = _manager()

    read = await _execute(manager, workspace, "read_file", {"path": str(skill_file)}, allowlist=allowlist)
    listed = await _execute(manager, workspace, "list_dir", {"path": str(skill_dir)}, allowlist=allowlist)
    written = await _execute(
        manager,
        workspace,
        "write_file",
        {"path": str(skill_file), "content": "replaced\n"},
        allowlist=allowlist,
    )
    edited = await _execute(
        manager,
        workspace,
        "edit_file",
        {"path": str(skill_file), "old_text": "replaced", "new_text": "edited"},
        allowlist=allowlist,
    )
    blocked = await _execute(
        manager,
        workspace,
        "write_file",
        {"path": str(tmp_path / "outside.txt"), "content": "no"},
        allowlist=allowlist,
    )

    assert "original" in read.content["content"]
    assert [entry["path"] for entry in listed.content["entries"]] == [str(skill_file)]
    assert written.content["success"] is True
    assert edited.content["success"] is True
    assert skill_file.read_text(encoding="utf-8") == "edited\n"
    assert blocked.content == {"error": "path is outside writable roots"}


def test_workspace_access_enforces_the_path_permission_matrix(tmp_path) -> None:
    workspace = tmp_path / "workspace"
    read_root = tmp_path / "read"
    write_root = tmp_path / "write"
    outside = tmp_path / "outside"
    for root in (workspace, read_root, write_root, outside):
        root.mkdir()
    access = WorkspaceAccess(
        workspace,
        extra_read_roots=(read_root,),
        extra_write_roots=(write_root,),
    )

    assert access.resolve_read("inside.txt") == workspace / "inside.txt"
    assert access.resolve_write("nested/../inside.txt") == workspace / "inside.txt"
    assert access.resolve_read(read_root / "read.txt") == read_root / "read.txt"
    assert access.resolve_write(read_root / "blocked.txt") == "path is outside writable roots"
    assert access.resolve_read(write_root / "shared.txt") == write_root / "shared.txt"
    assert access.resolve_write(write_root / "shared.txt") == write_root / "shared.txt"
    assert access.resolve_read(workspace / ".." / "outside" / "secret.txt") == (
        "path is outside readable roots"
    )

    escape = workspace / "escape"
    try:
        escape.symlink_to(outside, target_is_directory=True)
    except OSError as exc:
        pytest.skip(f"symlinks are not supported: {exc}")
    assert access.resolve_read(escape / "secret.txt") == "path is outside readable roots"
    assert access.resolve_write(escape / "created.txt") == "path is outside writable roots"


@pytest.mark.asyncio
async def test_read_only_allowlist_is_readable_but_not_writable(tmp_path) -> None:
    workspace = tmp_path / "workspace"
    read_root = tmp_path / "reference"
    workspace.mkdir()
    read_root.mkdir()
    target = read_root / "notes.txt"
    target.write_text("reference\n", encoding="utf-8")
    allowlist = PathAllowlist.from_roots(extra_read_roots=[read_root])
    manager = _manager()

    read, write = await asyncio.gather(
        _execute(
            manager,
            workspace,
            "read_file",
            {"path": str(target)},
            allowlist=allowlist,
        ),
        _execute(
            manager,
            workspace,
            "write_file",
            {"path": str(target), "content": "changed\n"},
            allowlist=allowlist,
        ),
    )

    assert "reference" in read.content["content"]
    assert write.content == {"error": "path is outside writable roots"}
    assert target.read_text(encoding="utf-8") == "reference\n"
