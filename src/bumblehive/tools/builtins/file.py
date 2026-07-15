import difflib
import os
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import fitz
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation as PptxPresentation

from ..adapters.function import CallableTool
from ..registry import ToolRegistry
from .state import BuiltinToolState, _session_scoped_file_handler
from .workspace import (
    _DEFAULT_IGNORE_DIRS,
    FileStates,
    WorkspaceAccess,
    current_file_states,
    current_workspace_access,
    is_binary_bytes,
)


_BLOCKED_DEVICE_PATHS = frozenset(
    {
        "/dev/console",
        "/dev/fd/0",
        "/dev/fd/1",
        "/dev/fd/2",
        "/dev/full",
        "/dev/random",
        "/dev/stderr",
        "/dev/stdin",
        "/dev/stdout",
        "/dev/tty",
        "/dev/urandom",
        "/dev/zero",
    }
)


def _is_blocked_device(path: str | Path) -> bool:
    raw = str(path)
    try:
        resolved = str(Path(raw).resolve())
    except (OSError, ValueError):
        resolved = raw

    if raw in _BLOCKED_DEVICE_PATHS or resolved in _BLOCKED_DEVICE_PATHS:
        return True
    if re.match(r"/proc/\d+/fd/[012]$", raw) or re.match(r"/proc/self/fd/[012]$", raw):
        return True
    if re.match(r"/proc/\d+/fd/[012]$", resolved) or re.match(r"/proc/self/fd/[012]$", resolved):
        return True
    return resolved.startswith("/dev/")


def _parse_page_range(pages: str, total_pages: int) -> tuple[int, int]:
    """Parse a 1-based page range into 0-based inclusive indexes."""
    parts = pages.strip().split("-")
    if not 1 <= len(parts) <= 2 or any(not part.strip() for part in parts):
        raise ValueError("invalid page range")

    start_page = int(parts[0])
    end_page = int(parts[-1])
    start = max(0, start_page - 1)
    end = min(end_page - 1, total_pages - 1)
    return start, end


def _collect_pptx_shape_text(shape: Any, out: list[str]) -> None:
    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        for sub_shape in sub_shapes:
            _collect_pptx_shape_text(sub_shape, out)
        return

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            row_text = "\t".join(cell for cell in cells if cell)
            if row_text:
                out.append(row_text)
        return

    text = getattr(shape, "text", "")
    if text:
        out.append(text)


_READ_FILE_DESCRIPTION = (
    "Read a UTF-8 text file or extract text from a PDF, DOCX, XLSX, or PPTX file. "
    "Text lines are returned as LINE_NUMBER| CONTENT. Use find_files or list_dir "
    "first when the path is uncertain, and read the relevant content before editing. "
    "Use offset and limit for large text files; output over 128,000 characters is "
    "truncated. Use force=true only to bypass unchanged-range deduplication."
)

_READ_FILE_PARAMETERS: dict[str, Any] = {
    "type": "object",
    "properties": {
        "path": {
            "type": "string",
            "description": (
                "File to read. Relative paths resolve from the workspace; absolute "
                "paths must be inside a readable root."
            ),
        },
        "offset": {
            "type": "integer",
            "description": "1-based line number to start reading from. Default 1.",
            "minimum": 1,
        },
        "limit": {
            "type": "integer",
            "description": "Maximum number of lines to read. Default 2000.",
            "minimum": 1,
        },
        "pages": {
            "type": "string",
            "description": (
                "PDF page range, for example '1-5'. Omit to read the first 20 pages; "
                "at most 20 pages are returned."
            ),
        },
        "force": {
            "type": "boolean",
            "description": (
                "Bypass unchanged-range deduplication and return content again. "
                "Default false."
            ),
        },
    },
    "required": ["path"],
    "additionalProperties": False,
}

_WRITE_FILE_DESCRIPTION = (
    "Create a UTF-8 text file or intentionally replace an entire file. Existing "
    "content is overwritten and missing parent directories are created. Prefer "
    "apply_patch for partial or multi-file changes, and edit_file for one narrow "
    "replacement."
)

_WRITE_FILE_PARAMETERS: dict[str, Any] = {
    "type": "object",
    "properties": {
        "path": {
            "type": "string",
            "description": (
                "File to create or replace. Relative paths resolve from the workspace; "
                "absolute paths must be inside a writable root."
            ),
        },
        "content": {
            "type": "string",
            "description": "Complete replacement content, up to 200,000 characters.",
        },
    },
    "required": ["path", "content"],
    "additionalProperties": False,
}

_LIST_DIR_DESCRIPTION = (
    "List the contents of a directory. Set recursive=true to explore nested "
    "structure. Common dependency, build, cache, and version-control directories "
    "are ignored."
)

_LIST_DIR_PARAMETERS: dict[str, Any] = {
    "type": "object",
    "properties": {
        "path": {
            "type": "string",
            "description": (
                "Directory to list. Relative paths resolve from the workspace; absolute "
                "paths must be inside a readable root."
            ),
        },
        "recursive": {
            "type": "boolean",
            "description": "Whether to recursively list nested contents. Default false.",
        },
        "max_entries": {
            "type": "integer",
            "description": "Maximum number of entries to return. Default 200.",
            "minimum": 1,
        },
    },
    "required": ["path"],
    "additionalProperties": False,
}

_EDIT_FILE_DESCRIPTION = (
    "Make a narrow text replacement in one UTF-8 file. Read the target first and "
    "copy old_text from current content. Use apply_patch for multi-file or larger "
    "changes. If old_text matches more than once, select a match with occurrence or "
    "line_hint, or set replace_all=true. Failed matches include diagnostics."
)

_EDIT_FILE_PARAMETERS: dict[str, Any] = {
    "type": "object",
    "properties": {
        "path": {
            "type": "string",
            "description": (
                "File to edit. Relative paths resolve from the workspace; absolute "
                "paths must be inside a writable root."
            ),
        },
        "old_text": {
            "type": "string",
            "description": "Current text to replace, preferably copied from read_file.",
        },
        "new_text": {
            "type": "string",
            "description": "Text that replaces the selected old_text match.",
        },
        "replace_all": {
            "type": "boolean",
            "description": "Replace every occurrence. Default false.",
        },
        "occurrence": {
            "type": "integer",
            "description": (
                "1-based occurrence to replace when old_text appears multiple times."
            ),
            "minimum": 1,
        },
        "line_hint": {
            "type": "integer",
            "description": (
                "Optional exact 1-based target line copied from read_file. "
                "The selected old_text match must cover this line."
            ),
            "minimum": 1,
        },
        "expected_replacements": {
            "type": "integer",
            "description": "Optional guard for the number of replacements.",
            "minimum": 1,
        },
    },
    "required": ["path", "old_text", "new_text"],
    "additionalProperties": False,
}


class WorkspaceFiles:
    _MAX_READ_CHARS = 128_000
    _DEFAULT_READ_LIMIT = 2000
    _MAX_PDF_PAGES = 20
    _MAX_WRITE_CHARS = 200_000
    _MAX_EDIT_CHARS = 1024 * 1024 * 1024
    _DEFAULT_LIST_ENTRIES = 200
    _MARKDOWN_EXTS = frozenset({".md", ".mdx", ".markdown"})
    _OFFICE_EXTS = frozenset({".docx", ".xlsx", ".pptx"})

    def __init__(
        self,
        file_states: FileStates | None = None,
    ) -> None:
        self._explicit_file_states = file_states
        self._fallback_file_states = FileStates()

    @property
    def file_states(self) -> FileStates:
        if self._explicit_file_states is not None:
            return self._explicit_file_states
        return current_file_states(self._fallback_file_states)

    def read_file(
        self,
        path: str,
        offset: int = 1,
        limit: int | None = None,
        pages: str | None = None,
        force: bool = False,
    ) -> dict[str, Any]:
        if _is_blocked_device(path):
            return {"error": "reading device paths is blocked", "path": path}
        access = self._access()
        resolved = access.resolve_read(path)
        if isinstance(resolved, str):
            return {"error": resolved}
        if _is_blocked_device(resolved):
            return {"error": "reading device paths is blocked", "path": str(resolved)}
        if not resolved.exists() or not resolved.is_file():
            return {"error": "path does not exist or is not a file", "path": str(resolved)}
        read_offset = offset
        read_limit = limit
        read_offset = max(read_offset, 1)

        suffix = resolved.suffix.lower()
        if suffix == ".pdf":
            return self._read_pdf(resolved, pages)
        if suffix in self._OFFICE_EXTS:
            return self._read_office_doc(resolved)

        try:
            raw = resolved.read_bytes()
        except UnicodeDecodeError:
            return {"error": "file is not valid UTF-8 text", "path": str(resolved)}
        except OSError as exc:
            return {"error": str(exc), "path": str(resolved)}

        if not raw:
            return {
                "path": str(resolved),
                "start_line": 1,
                "end_line": 0,
                "total_lines": 0,
                "content": "",
                "truncated": False,
                "deduplicated": False,
            }
        if is_binary_bytes(raw):
            return {"error": "file is not valid UTF-8 text", "path": str(resolved)}

        if not force and self.file_states.is_unchanged(
            resolved,
            offset=read_offset,
            limit=read_limit,
        ):
            return {
                "path": str(resolved),
                "content": "",
                "deduplicated": True,
                "message": "file range already read; pass force=true to read again",
                "offset": read_offset,
                "limit": read_limit,
            }

        try:
            text_content = raw.decode("utf-8").replace("\r\n", "\n")
        except UnicodeDecodeError:
            return {"error": "file is not valid UTF-8 text", "path": str(resolved)}

        lines = text_content.splitlines()
        total_lines = len(lines)
        if read_offset > total_lines and total_lines > 0:
            return {
                "error": f"offset {read_offset} is beyond end of file",
                "path": str(resolved),
                "total_lines": total_lines,
            }

        start = read_offset - 1
        effective_limit = read_limit or self._DEFAULT_READ_LIMIT
        end = min(start + effective_limit, total_lines)
        selected_lines = lines[start:end]
        numbered = [
            f"{line_number}| {line}"
            for line_number, line in enumerate(selected_lines, start=start + 1)
        ]

        truncated = False
        content = "\n".join(numbered)
        if len(content) > self._MAX_READ_CHARS:
            truncated = True
            trimmed: list[str] = []
            char_count = 0
            for line in numbered:
                next_count = char_count + len(line) + (1 if trimmed else 0)
                if next_count > self._MAX_READ_CHARS:
                    if not trimmed:
                        trimmed.append(line[:self._MAX_READ_CHARS])
                    break
                trimmed.append(line)
                char_count = next_count
            content = "\n".join(trimmed)
            end = start + len(trimmed)
        has_more = end < total_lines

        self.file_states.record_read(resolved, offset=read_offset, limit=read_limit)
        return {
            "path": str(resolved),
            "start_line": start + 1,
            "end_line": end,
            "total_lines": total_lines,
            "content": content,
            "truncated": truncated,
            "has_more": has_more,
            "deduplicated": False,
            "next_offset": end + 1 if has_more else None,
        }

    def _read_pdf(self, path: Path, pages: str | None) -> dict[str, Any]:
        try:
            doc = fitz.open(str(path))
        except Exception as exc:
            return {"error": f"error reading PDF: {exc}", "path": str(path)}

        try:
            total_pages = len(doc)
            if total_pages == 0:
                self.file_states.record_read(path)
                return {
                    "path": str(path),
                    "file_type": "pdf",
                    "content": "",
                    "pages": None,
                    "total_pages": 0,
                    "truncated": False,
                    "has_more": False,
                }

            if pages:
                try:
                    start, end = _parse_page_range(pages, total_pages)
                except ValueError:
                    return {
                        "error": f"invalid page range '{pages}'. Use format like '1-5'.",
                        "path": str(path),
                        "total_pages": total_pages,
                    }
                if start > end or start >= total_pages:
                    return {
                        "error": f"page range '{pages}' is out of bounds",
                        "path": str(path),
                        "total_pages": total_pages,
                    }
            else:
                start = 0
                end = min(total_pages - 1, self._MAX_PDF_PAGES - 1)

            page_count = end - start + 1
            if page_count > self._MAX_PDF_PAGES:
                end = start + self._MAX_PDF_PAGES - 1

            parts: list[str] = []
            for index in range(start, end + 1):
                text = doc[index].get_text().strip()
                if text:
                    parts.append(f"--- Page {index + 1} ---\n{text}")

            content = "\n\n".join(parts)
            content, truncated = self._truncate_read_content(content)
            has_more = end < total_pages - 1
            self.file_states.record_read(path)
            return {
                "path": str(path),
                "file_type": "pdf",
                "content": content,
                "pages": f"{start + 1}-{end + 1}",
                "total_pages": total_pages,
                "truncated": truncated,
                "has_more": has_more,
                "next_pages": (
                    f"{end + 2}-{min(end + 1 + self._MAX_PDF_PAGES, total_pages)}"
                    if has_more
                    else None
                ),
                "message": None if content else "PDF has no extractable text",
            }
        finally:
            doc.close()

    def _read_office_doc(self, path: Path) -> dict[str, Any]:
        suffix = path.suffix.lower()
        try:
            if suffix == ".docx":
                content = self._extract_docx(path)
            elif suffix == ".xlsx":
                content = self._extract_xlsx(path)
            elif suffix == ".pptx":
                content = self._extract_pptx(path)
            else:
                return {"error": f"unsupported document format: {suffix}", "path": str(path)}
        except Exception as exc:
            return {"error": f"error reading {suffix.upper()} file: {exc}", "path": str(path)}

        content, truncated = self._truncate_read_content(content)
        self.file_states.record_read(path)
        return {
            "path": str(path),
            "file_type": suffix.lstrip("."),
            "content": content,
            "truncated": truncated,
            "message": None if content else f"{suffix.upper().lstrip('.')} has no extractable text",
        }

    def _extract_docx(self, path: Path) -> str:
        doc = DocxDocument(path)
        parts: list[str] = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                row_text = "\t".join(cell for cell in cells if cell)
                if row_text:
                    parts.append(row_text)
        return "\n\n".join(parts)

    def _extract_xlsx(self, path: Path) -> str:
        workbook = load_workbook(path, read_only=True, data_only=True)
        try:
            sheets: list[str] = []
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                rows: list[str] = []
                for row in worksheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        rows.append(row_text)
                if rows:
                    sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
            return "\n\n".join(sheets)
        finally:
            workbook.close()

    def _extract_pptx(self, path: Path) -> str:
        presentation = PptxPresentation(path)
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, 1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                _collect_pptx_shape_text(shape, slide_text)
            if slide_text:
                slides.append(f"--- Slide {index} ---\n" + "\n".join(slide_text))
        return "\n\n".join(slides)

    def _truncate_read_content(self, content: str) -> tuple[str, bool]:
        if len(content) <= self._MAX_READ_CHARS:
            return content, False
        return content[:self._MAX_READ_CHARS], True

    def write_file(self, path: str, content: str) -> dict[str, Any]:
        access = self._access()
        resolved = access.resolve_write(path)
        if isinstance(resolved, str):
            return {"error": resolved}
        if len(content) > self._MAX_WRITE_CHARS:
            return {"error": f"content is too large; max {self._MAX_WRITE_CHARS} chars"}

        resolved.parent.mkdir(parents=True, exist_ok=True)
        resolved.write_text(content, encoding="utf-8")
        self.file_states.record_write(resolved)
        return {
            "path": str(resolved),
            "bytes_written": len(content.encode("utf-8")),
            "success": True,
        }

    def list_dir(
        self,
        path: str,
        recursive: bool = False,
        max_entries: int | None = None,
    ) -> dict[str, Any]:
        access = self._access()
        resolved = access.resolve_read(path)
        if isinstance(resolved, str):
            return {"error": resolved}
        if not resolved.exists():
            return {"error": "directory does not exist", "path": str(resolved)}
        if not resolved.is_dir():
            return {"error": "path is not a directory", "path": str(resolved)}

        cap = max_entries or self._DEFAULT_LIST_ENTRIES
        entries: list[dict[str, str]] = []
        total_entries = 0

        iterator = self._iter_dir(resolved, recursive=recursive)
        for item in iterator:
            total_entries += 1
            if len(entries) >= cap:
                continue

            entries.append(
                {
                    "path": access.relative_display_path(item),
                    "type": "directory" if item.is_dir() else "file",
                }
            )

        return {
            "path": str(resolved),
            "entries": entries,
            "total_entries": total_entries,
            "truncated": total_entries > cap,
        }

    @staticmethod
    def _iter_dir(path: Path, *, recursive: bool) -> list[Path]:
        if not recursive:
            return sorted(
                item for item in path.iterdir()
                if item.name not in _DEFAULT_IGNORE_DIRS
            )

        entries: list[Path] = []
        for dirpath, dirnames, filenames in os.walk(path):
            dirnames[:] = sorted(name for name in dirnames if name not in _DEFAULT_IGNORE_DIRS)
            current = Path(dirpath)
            entries.extend(current / dirname for dirname in dirnames)
            entries.extend(current / filename for filename in sorted(filenames))
        return entries

    def edit_file(
        self,
        path: str,
        old_text: str,
        new_text: str,
        replace_all: bool = False,
        occurrence: int | None = None,
        line_hint: int | None = None,
        expected_replacements: int | None = None,
    ) -> dict[str, Any]:
        access = self._access()
        resolved = access.resolve_write(path)
        if isinstance(resolved, str):
            return {"error": resolved}
        if replace_all and occurrence is not None:
            return {"error": "occurrence cannot be used with replace_all=true"}
        if replace_all and line_hint is not None:
            return {"error": "line_hint cannot be used with replace_all=true"}
        if occurrence is not None and line_hint is not None:
            return {"error": "line_hint cannot be used with occurrence"}

        if not resolved.exists():
            if old_text == "":
                resolved.parent.mkdir(parents=True, exist_ok=True)
                resolved.write_text(new_text, encoding="utf-8")
                self.file_states.record_write(resolved)
                return {"path": str(resolved), "created": True, "success": True}
            return {"error": "path does not exist or is not a file", "path": str(resolved)}
        if not resolved.is_file():
            return {"error": "path is not a file", "path": str(resolved)}
        if resolved.stat().st_size > self._MAX_EDIT_CHARS:
            return {"error": f"file is too large to edit; max {self._MAX_EDIT_CHARS} bytes"}
        warning = self.file_states.check_read(resolved) if old_text != "" else None

        try:
            raw = resolved.read_bytes()
            content = raw.decode("utf-8").replace("\r\n", "\n")
        except UnicodeDecodeError:
            return {"error": "file is not valid UTF-8 text", "path": str(resolved)}

        if old_text == "":
            if content.strip():
                return {"error": "cannot create file because it already exists and is not empty"}
            resolved.write_text(new_text, encoding="utf-8")
            self.file_states.record_write(resolved)
            return {"path": str(resolved), "replacements": 1, "success": True}

        norm_old = old_text.replace("\r\n", "\n")
        norm_new = new_text.replace("\r\n", "\n")
        if resolved.suffix.lower() not in self._MARKDOWN_EXTS:
            norm_new = _strip_trailing_ws(norm_new)

        matches = _find_matches(content, norm_old)
        if not matches:
            return {
                "error": "old_text not found",
                "path": str(resolved),
                "nearest_match": _nearest_match(content, norm_old),
            }

        if replace_all:
            selected = matches
        elif occurrence is not None:
            if occurrence > len(matches):
                return {
                    "error": "occurrence is out of range",
                    "occurrences": len(matches),
                }
            selected = [matches[occurrence - 1]]
        elif line_hint is not None:
            candidates = [
                match for match in matches
                if _match_covers_line(match, line_hint)
            ]
            if not candidates:
                return {
                    "error": "line_hint does not match the old_text location",
                    "line_hint": line_hint,
                    "occurrences": len(matches),
                    "match_locations": [
                        _match_location(match) for match in matches[:5]
                    ],
                }
            if len(candidates) > 1:
                return {
                    "error": (
                        "line_hint is ambiguous; old_text appears multiple times "
                        "on the target line"
                    ),
                    "line_hint": line_hint,
                    "occurrences": len(candidates),
                    "match_locations": [
                        _match_location(match) for match in candidates[:5]
                    ],
                }
            selected = candidates
        else:
            if len(matches) > 1:
                return {
                    "error": "old_text appears multiple times; provide occurrence or replace_all",
                    "occurrences": len(matches),
                    "lines": [match.line for match in matches],
                }
            selected = [matches[0]]

        if (
            expected_replacements is not None
            and len(selected) != expected_replacements
        ):
            return {
                "error": "replacement count did not match expected_replacements",
                "expected_replacements": expected_replacements,
                "actual_replacements": len(selected),
            }

        updated = content
        for match in reversed(selected):
            replacement = _preserve_quote_style(norm_old, match.text, norm_new)
            replacement = _reindent_like_match(norm_old, match.text, replacement)
            end = match.end
            if replacement == "" and not match.text.endswith("\n") and updated[end:end + 1] == "\n":
                end += 1
            updated = updated[:match.start] + replacement + updated[end:]

        if b"\r\n" in raw:
            updated = updated.replace("\n", "\r\n")
        resolved.write_bytes(updated.encode("utf-8"))
        self.file_states.record_write(resolved)
        result = {
            "path": str(resolved),
            "replacements": len(selected),
            "success": True,
        }
        if warning:
            result["warning"] = warning
        return result

    def _access(self) -> WorkspaceAccess:
        return current_workspace_access()


@dataclass(frozen=True)
class MatchSpan:
    start: int
    end: int
    text: str
    line: int


def _strip_trailing_ws(text: str) -> str:
    return "\n".join(line.rstrip() for line in text.split("\n"))


def _match_end_line(match: MatchSpan) -> int:
    comparable = match.text[:-1] if match.text.endswith("\n") else match.text
    return match.line + comparable.count("\n")


def _match_covers_line(match: MatchSpan, line: int) -> bool:
    return match.line <= line <= _match_end_line(match)


def _match_location(match: MatchSpan) -> dict[str, int]:
    return {
        "start_line": match.line,
        "end_line": _match_end_line(match),
    }


def _find_exact_matches(content: str, old_text: str) -> list[MatchSpan]:
    matches: list[MatchSpan] = []
    start = 0
    while True:
        index = content.find(old_text, start)
        if index == -1:
            return matches
        matches.append(
            MatchSpan(
                start=index,
                end=index + len(old_text),
                text=content[index:index + len(old_text)],
                line=content.count("\n", 0, index) + 1,
            )
        )
        start = index + max(1, len(old_text))


def _find_trim_matches(
    content: str,
    old_text: str,
    *,
    normalize_quotes: bool = False,
) -> list[MatchSpan]:
    old_lines = old_text.splitlines()
    if not old_lines:
        return []

    content_lines = content.splitlines()
    keepends = content.splitlines(keepends=True)
    if len(content_lines) < len(old_lines):
        return []

    offsets: list[int] = []
    pos = 0
    for line in keepends:
        offsets.append(pos)
        pos += len(line)
    offsets.append(pos)

    if normalize_quotes:
        stripped_old = [_normalize_quotes(line.strip()) for line in old_lines]
    else:
        stripped_old = [line.strip() for line in old_lines]

    matches: list[MatchSpan] = []
    window_size = len(stripped_old)
    for index in range(len(content_lines) - window_size + 1):
        window = content_lines[index:index + window_size]
        if normalize_quotes:
            comparable = [_normalize_quotes(line.strip()) for line in window]
        else:
            comparable = [line.strip() for line in window]
        if comparable != stripped_old:
            continue
        start = offsets[index]
        end = offsets[index + window_size]
        if keepends[index + window_size - 1].endswith("\n"):
            end -= 1
        matches.append(
            MatchSpan(
                start=start,
                end=end,
                text=content[start:end],
                line=index + 1,
            )
        )
    return matches


def _find_quote_matches(content: str, old_text: str) -> list[MatchSpan]:
    norm_content = _normalize_quotes(content)
    norm_old = _normalize_quotes(old_text)
    matches: list[MatchSpan] = []
    start = 0
    while True:
        index = norm_content.find(norm_old, start)
        if index == -1:
            return matches
        matches.append(
            MatchSpan(
                start=index,
                end=index + len(old_text),
                text=content[index:index + len(old_text)],
                line=content.count("\n", 0, index) + 1,
            )
        )
        start = index + max(1, len(norm_old))


def _find_matches(content: str, old_text: str) -> list[MatchSpan]:
    for matcher in (
        lambda: _find_exact_matches(content, old_text),
        lambda: _find_trim_matches(content, old_text),
        lambda: _find_trim_matches(content, old_text, normalize_quotes=True),
        lambda: _find_quote_matches(content, old_text),
    ):
        matches = matcher()
        if matches:
            return matches
    return []


def _normalize_quotes(text: str) -> str:
    return text.translate(
        str.maketrans(
            {
                "\u2018": "'",
                "\u2019": "'",
                "\u201c": '"',
                "\u201d": '"',
            }
        )
    )


def _leading_ws(line: str) -> str:
    return line[: len(line) - len(line.lstrip(" \t"))]


def _reindent_like_match(old_text: str, actual_text: str, new_text: str) -> str:
    old_lines = old_text.split("\n")
    actual_lines = actual_text.split("\n")
    if len(old_lines) != len(actual_lines):
        return new_text

    comparable = [
        (old_line, actual_line)
        for old_line, actual_line in zip(old_lines, actual_lines)
        if old_line.strip() and actual_line.strip()
    ]
    if not comparable:
        return new_text
    if any(
        _normalize_quotes(old_line.strip()) != _normalize_quotes(actual_line.strip())
        for old_line, actual_line in comparable
    ):
        return new_text

    old_ws = _leading_ws(comparable[0][0])
    actual_ws = _leading_ws(comparable[0][1])
    if actual_ws == old_ws:
        return new_text
    if old_ws:
        if not actual_ws.startswith(old_ws):
            return new_text
        delta = actual_ws[len(old_ws):]
    else:
        delta = actual_ws
    if not delta:
        return new_text
    return "\n".join((delta + line) if line else line for line in new_text.split("\n"))


def _preserve_quote_style(old_text: str, actual_text: str, new_text: str) -> str:
    if _normalize_quotes(old_text.strip()) != _normalize_quotes(actual_text.strip()):
        return new_text
    if old_text == actual_text:
        return new_text

    styled = new_text
    if any(ch in actual_text for ch in ("\u201c", "\u201d")) and '"' in styled:
        styled = _curly_double_quotes(styled)
    if any(ch in actual_text for ch in ("\u2018", "\u2019")) and "'" in styled:
        styled = _curly_single_quotes(styled)
    return styled


def _curly_double_quotes(text: str) -> str:
    parts: list[str] = []
    opening = True
    for char in text:
        if char == '"':
            parts.append("\u201c" if opening else "\u201d")
            opening = not opening
        else:
            parts.append(char)
    return "".join(parts)


def _curly_single_quotes(text: str) -> str:
    parts: list[str] = []
    opening = True
    for index, char in enumerate(text):
        if char != "'":
            parts.append(char)
            continue
        prev_char = text[index - 1] if index > 0 else ""
        next_char = text[index + 1] if index + 1 < len(text) else ""
        if prev_char.isalnum() and next_char.isalnum():
            parts.append("\u2019")
            continue
        parts.append("\u2018" if opening else "\u2019")
        opening = not opening
    return "".join(parts)


def _nearest_match(content: str, old_text: str) -> dict[str, Any] | None:
    old_lines = old_text.splitlines() or [old_text]
    content_lines = content.splitlines()
    if not content_lines:
        return None

    window = max(1, len(old_lines))
    best_ratio = -1.0
    best_index = 0
    best_lines: list[str] = []

    for index in range(max(1, len(content_lines) - window + 1)):
        candidate = content_lines[index:index + window]
        ratio = difflib.SequenceMatcher(
            None,
            "\n".join(old_lines),
            "\n".join(candidate),
        ).ratio()
        if ratio > best_ratio:
            best_ratio = ratio
            best_index = index
            best_lines = candidate

    return {
        "line": best_index + 1,
        "similarity": round(best_ratio, 3),
        "content": "\n".join(best_lines),
    }


def _workspace_files_from_state(state: BuiltinToolState) -> WorkspaceFiles:
    files = state.workspace_files
    if not isinstance(files, WorkspaceFiles):
        files = WorkspaceFiles()
        state.workspace_files = files
    return files


def register_read_file_tool(
    registry: ToolRegistry,
    *,
    state: BuiltinToolState,
) -> CallableTool:
    """Register the read_file tool on a registry."""
    files = _workspace_files_from_state(state)
    return registry.register(
        CallableTool(
            name="read_file",
            description=_READ_FILE_DESCRIPTION,
            parameters=_READ_FILE_PARAMETERS,
            handler=_session_scoped_file_handler(state, files.read_file),
            read_only=True,
        )
    )


def register_write_file_tool(
    registry: ToolRegistry,
    *,
    state: BuiltinToolState,
) -> CallableTool:
    """Register the write_file tool on a registry."""
    files = _workspace_files_from_state(state)
    return registry.register(
        CallableTool(
            name="write_file",
            description=_WRITE_FILE_DESCRIPTION,
            parameters=_WRITE_FILE_PARAMETERS,
            handler=_session_scoped_file_handler(state, files.write_file),
            exclusive=True,
        )
    )


def register_list_dir_tool(
    registry: ToolRegistry,
    *,
    state: BuiltinToolState,
) -> CallableTool:
    """Register the list_dir tool on a registry."""
    files = _workspace_files_from_state(state)
    return registry.register(
        CallableTool(
            name="list_dir",
            description=_LIST_DIR_DESCRIPTION,
            parameters=_LIST_DIR_PARAMETERS,
            handler=files.list_dir,
            read_only=True,
        )
    )


def register_edit_file_tool(
    registry: ToolRegistry,
    *,
    state: BuiltinToolState,
) -> CallableTool:
    """Register the edit_file tool on a registry."""
    files = _workspace_files_from_state(state)
    return registry.register(
        CallableTool(
            name="edit_file",
            description=_EDIT_FILE_DESCRIPTION,
            parameters=_EDIT_FILE_PARAMETERS,
            handler=_session_scoped_file_handler(state, files.edit_file),
            exclusive=True,
        )
    )
